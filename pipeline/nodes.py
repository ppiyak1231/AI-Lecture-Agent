import os
import json
from pathlib import Path
from typing import List, Dict, TypedDict, Any

import streamlit as st
from pptx import Presentation
import openai
from openai import OpenAI
from langchain_openai import ChatOpenAI
from langchain_core.messages import SystemMessage, HumanMessage

# 미디어 및 챗봇 헬퍼 임포트
from utils.media_helper import clean_text, export_slide_as_png, render_mp4, concat_videos_ffmpeg

# LLM & TTS 모델 상수
LLM_MODEL = "gpt-4o-mini"
TTS_MODEL = "tts-1"

# -----------------------------------------------------------------------------
# 1. 1번: 강의 영상 제작용 State & Nodes
# -----------------------------------------------------------------------------
class State(TypedDict, total=False):
    pptx_path: str
    work_dir: str
    prompt: Dict            # voice, tone, style, persona, level
    openai_key: str
    tavily_key: str
    
    slides: List[Dict]      # title, texts, tables, snap
    n_slides: int
    slide_index: int
    
    page_content: str
    script: str
    audio: str
    video_path: str
    
    all_contents: List[str]
    all_scripts: List[str]
    all_audios: List[str]
    all_videos: List[str]
    
    final_video: str
    next_worker: str        # Supervisor가 지명한 다음 작업자 (Parser, Summarizer, ScriptWriter, MediaCreator, Accumulator, ConcatVideo, FINISH)

def node_lecture_supervisor(state: State) -> State:
    """
    [Lecture Supervisor Agent]
    현재 제작 현황 상태를 보고받고 지능적으로 다음 단계의 Worker를 배정합니다.
    저렴한 gpt-4o-mini 모델을 사용하여 제어를 수행합니다.
    """
    st.session_state["progress_status"] = "🧠 [Lecture Supervisor] 상태를 모니터링하여 적합한 Worker를 배치 중..."
    
    # 만약 슬라이드 정보가 파싱되지 않았다면 우선 파싱 지명
    if "slides" not in state or not state["slides"]:
        state["next_worker"] = "Parser"
        return state
        
    idx = state.get("slide_index", 0)
    n_slides = state.get("n_slides", 0)
    
    # 모든 슬라이드 공정이 끝났다면 병합 지명
    if idx >= n_slides:
        if "final_video" not in state or not state["final_video"]:
            state["next_worker"] = "ConcatVideo"
        else:
            state["next_worker"] = "FINISH"
        return state
        
    # 현재 슬라이드의 제작 현황을 분석
    has_summary = len(state.get("all_contents", [])) > idx
    has_script = len(state.get("all_scripts", [])) > idx
    has_media = len(state.get("all_videos", [])) > idx
    
    # 순서대로 요약 -> 대본 -> 오디오/비디오 제작 -> 누적(Accumulate) 지명
    if not has_summary:
        state["next_worker"] = "Summarizer"
    elif not has_script:
        state["next_worker"] = "ScriptWriter"
    elif not has_media:
        state["next_worker"] = "MediaCreator"
    else:
        state["next_worker"] = "Accumulator"
        
    return state

def node_parse_all(state: State) -> State:
    """[Parser Worker] PPTX 슬라이드를 읽어 구조화된 텍스트와 레이아웃을 정밀 파싱합니다."""
    st.session_state["progress_status"] = "📄 [Parser Worker] 프레젠테이션 분석 및 이미지 스냅샷 생성 중..."
    pres = Presentation(state["pptx_path"])
    
    slides_out = []
    for idx, slide in enumerate(pres.slides):
        texts, tables = [], []
        title = ""
        
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text
            
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = clean_text(para.text)
                    if line: 
                        texts.append(line)
            elif shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    rows.append([clean_text(cell.text) for cell in row.cells])
                tables.append(rows)
                
        # 슬라이드 렌더링 이미지 준비 (soffice 혹은 Pillow fallback 자동 적용)
        snap = export_slide_as_png(state, idx, title, texts)
        
        slides_out.append({
            "index": idx + 1,
            "title": title,
            "texts": texts,
            "tables": tables,
            "snap": snap
        })
        
    state["slides"] = slides_out
    state["n_slides"] = len(slides_out)
    state["slide_index"] = 0
    state["all_contents"] = []
    state["all_scripts"] = []
    state["all_audios"] = []
    state["all_videos"] = []
    return state

def node_generate_text(state: State) -> State:
    """[Summarizer Worker] 슬라이드의 핵심 정보를 가치가 높은 명확한 문장으로 요약합니다."""
    idx = state["slide_index"]
    st.session_state["progress_status"] = f"✍️ [Summarizer Worker] {idx+1}/{state['n_slides']} 슬라이드 요약 작성 중..."
    
    cur_slide = state["slides"][idx]
    texts = cur_slide.get("texts", [])
    tables = cur_slide.get("tables", [])
    
    table_snip = ""
    if tables:
        try:
            table_snip = "\n".join([" | ".join(map(str, r)) for r in tables[0][:5]])
        except Exception:
            table_snip = str(tables[0][:5])
            
    prompt_style = state["prompt"].get("style", "핵심 요점 중심")
    
    sys_msg = """
    역할: 당신은 발표 슬라이드의 정보(텍스트, 표)를 상세히 핵심만 요약하는 전문 AI 강사 조교입니다.
    규칙:
    - 4~6문장으로 구체적이고 체계적인 정보를 정리하세요.
    - 불릿포인트(*, -)나 번호 매기기를 절대 사용하지 말고, 하나의 완성된 자연스러운 단락으로 작성하세요.
    - 표나 구조화된 정보의 경우 발표에 유용한 가치를 뽑아서 텍스트로 녹이세요.
    """
    
    user_text = (
        f"[슬라이드 제목]\n{cur_slide.get('title', '없음')}\n\n"
        f"[텍스트]\n{texts if texts else '(텍스트 없음)'}\n\n"
        f"[표 요약]\n{table_snip if table_snip else '(표 없음)'}\n\n"
        f"[요약 가이드라인]\n{prompt_style}"
    )
    
    chat = ChatOpenAI(model=LLM_MODEL, temperature=0.3, openai_api_key=state["openai_key"])
    response = chat.invoke([
        SystemMessage(content=sys_msg),
        HumanMessage(content=user_text)
    ])
    
    state["page_content"] = response.content
    return state

def node_generate_script(state: State) -> State:
    """[ScriptWriter Worker] 요약된 강의 지식을 자연스러운 구어체 대본으로 풀어냅니다."""
    idx = state["slide_index"]
    st.session_state["progress_status"] = f"📝 [ScriptWriter Worker] {idx+1}/{state['n_slides']} 슬라이드 강의 대본 작성 중..."
    
    page_content = state["page_content"]
    prompt = state["prompt"]
    persona = prompt.get("persona", "친절한 AI 강사")
    tone = prompt.get("tone", "친절하고 부드럽게, 명료하게 읽어주세요.")
    level = prompt.get("level", "비전공자")
    
    sys_msg = f"""
    역할: 당신은 {persona} 스타일의 자연스러운 강의 발표 대본 전문가입니다.
    목표: 제공받은 [슬라이드 요약 정보]를 토대로 청중({level} 대상)이 귀에 쏙쏙 들어오도록 말하는 구어체 대본을 작성하세요.
    
    규칙:
    - 발표 시간은 약 40~60초 분량으로 정갈하게 작성하세요.
    - 서론(도입) -> 본론(설명) -> 결론(요약/브릿지)의 깔끔한 흐름을 지키세요.
    - 반드시 실제 입으로 소리 내어 말하는 부드러운 한국어 구어체 (~해요, ~합니다)로만 스크립트를 반환하세요.
    - 슬라이드 정보 이외의 임의의 과장되거나 허위의 정보를 창작하지 마세요.
    """
    
    user_msg = f"[슬라이드 요약 정보]\n{page_content}\n\n[말조 및 톤 가이드라인]\n{tone}"
    
    chat = ChatOpenAI(model=LLM_MODEL, temperature=0.5, openai_api_key=state["openai_key"])
    script = chat.invoke([
        SystemMessage(content=sys_msg),
        HumanMessage(content=user_msg)
    ]).content
    
    state["script"] = script
    return state

def node_generate_media(state: State) -> State:
    """[MediaCreator Worker] OpenAI TTS와 FFmpeg을 활용해 음성을 합성하고 슬라이드 비디오를 제작합니다."""
    idx = state["slide_index"]
    st.session_state["progress_status"] = f"🎬 [MediaCreator Worker] {idx+1}/{state['n_slides']} 슬라이드 오디오/비디오 조각 빌드 중..."
    
    # 1. TTS 생성
    script = state["script"]
    voice = "alloy"
    persona = state["prompt"].get("persona", "친절한 AI 강사")
    if "현업 전문가" in persona:
        voice = "onyx"
    elif "교수" in persona:
        voice = "nova"
    elif "선배" in persona:
        voice = "shimmer"
        
    try:
        client = OpenAI(api_key=state["openai_key"])
        response = client.audio.speech.create(
            model=TTS_MODEL,
            voice=voice,
            input=script
        )
        
        mp3_path = os.path.join(state["work_dir"], f"narration_{idx}.mp3")
        with open(mp3_path, "wb") as f:
            f.write(response.content)
            
        state["audio"] = mp3_path
    except Exception as e:
        raise Exception(f"OpenAI TTS 음성 생성 실패: {str(e)}")
        
    # 2. 비디오 합성
    slide_image = state["slides"][idx]["snap"]
    audio_path = state["audio"]
    work_dir = state["work_dir"]
    
    out_path = os.path.join(work_dir, f"slide_{idx}_lecture.mp4")
    success = render_mp4(slide_image, audio_path, out_path)
    
    if not success:
        raise Exception(f"{idx+1}번째 슬라이드 동영상 생성(FFmpeg)에 실패했습니다.")
        
    state["video_path"] = out_path
    return state

def node_accumulate(state: State) -> State:
    """[Accumulator Worker] 생성된 슬라이드별 결과 데이터셋을 메인 전역 상태에 누적 기록합니다."""
    state["all_contents"].append(state["page_content"])
    state["all_scripts"].append(state["script"])
    state["all_audios"].append(state["audio"])
    state["all_videos"].append(state["video_path"])
    state["slide_index"] += 1
    return state

def node_concat_video(state: State) -> State:
    """[ConcatVideo Worker] 모든 동영상 파트를 결합하여 최종 명품 영상물로 내보냅니다."""
    st.session_state["progress_status"] = "🔗 [ConcatVideo Worker] 슬라이드별 비디오 조각들을 모아 최종 병합 영상 렌더링 중..."
    video_paths = [v for v in state["all_videos"] if v and os.path.exists(v)]
    
    if not video_paths:
        raise ValueError("합칠 동영상 목록이 비어 있습니다.")
        
    out_path = os.path.join(state["work_dir"], "final_lecture.mp4")
    concat_videos_ffmpeg(video_paths, out_path)
    
    state["final_video"] = out_path
    return state


# -----------------------------------------------------------------------------
# 2. 2번: RAG 챗봇용 State & Nodes
# -----------------------------------------------------------------------------
class ChatState(TypedDict, total=False):
    user_query: str
    chat_history: List[Dict[str, str]]
    openai_key: str
    tavily_key: str
    vector_store: Any
    
    retrieved_context: str
    web_search_context: str
    
    next_worker: str        # PPT_Retriever, Web_Searcher, Response_Generator, FINISH
    final_response: str

def node_chat_supervisor(state: ChatState) -> ChatState:
    """
    [Chat Supervisor Agent]
    사용자의 질문 의도 및 수집된 컨텍스트 정보를 분석하여 지능적인 작업 라우팅을 수행합니다.
    비용 효율이 탁월한 gpt-4o-mini 모델을 사용하여 지능적으로 감독 제어합니다.
    """
    # 1단계: RAG 검색이 우선 수행되지 않았다면 RAG Retriever 호출 지명
    if "retrieved_context" not in state or state["retrieved_context"] is None:
        state["next_worker"] = "PPT_Retriever"
        return state
        
    # 2단계: RAG 결과에 대한 요약을 평가하여 외부 지식 검색(Tavily)이 필요한지 LLM이 판단
    retrieved = state.get("retrieved_context", "")
    query = state.get("user_query", "")
    
    # Tavily Key가 있고 아직 웹 서치가 미수행되었을 때, 웹서치가 필요한지 여부 판단
    if state.get("tavily_key") and ("web_search_context" not in state or state["web_search_context"] is None):
        sys_msg = """
        역할: 당신은 사용자의 질문이 제공받은 PPT 검색 결과(Retrieved Context) 내에서 완벽하게 대답할 수 있는지 확인하는 감독관 에이전트입니다.
        규칙:
        - PPT 내부 검색 결과만으로 답변하기에 정보가 부족하거나, 전혀 관계가 없는 외부 지식을 묻는 질문인 경우 'YES'를 반환하세요.
        - 충분히 PPT 내부 내용만으로 풍성하게 대답할 수 있다면 'NO'를 반환하세요.
        - 오직 단 두 글자 ('YES' 또는 'NO')로만 답변하세요.
        """
        user_msg = f"[PPT 검색 결과]\n{retrieved}\n\n[사용자 질문]\n{query}"
        
        chat = ChatOpenAI(model=LLM_MODEL, temperature=0.1, openai_api_key=state["openai_key"])
        res = chat.invoke([
            SystemMessage(content=sys_msg),
            HumanMessage(content=user_msg)
        ]).content.strip().upper()
        
        if "YES" in res:
            state["next_worker"] = "Web_Searcher"
            return state
            
    # 정보가 수집되었거나 웹 검색이 불가할 때 -> 최종 답변 성형 지명
    if "final_response" not in state or not state["final_response"]:
        state["next_worker"] = "Response_Generator"
    else:
        state["next_worker"] = "FINISH"
        
    return state

def node_ppt_retriever(state: ChatState) -> ChatState:
    """[PPT Retriever Worker] FAISS 로컬 벡터스토어에서 질문에 매칭되는 슬라이드 내용을 검색합니다."""
    query = state["user_query"]
    vector_store = state["vector_store"]
    
    # 상위 2개의 연관성 높은 정보 획득
    search_results = vector_store.similarity_search_with_score(query, k=2)
    
    context_docs = []
    min_score = 999.0
    for doc, score in search_results:
        context_docs.append(doc.page_content)
        if score < min_score:
            min_score = score
            
    # 유사도 점수가 낮아(L2 distance 기준 0.95 초과) 신뢰할 수 없는 경우, PPT에 아예 없는 것으로 간주
    if min_score >= 0.95:
        state["retrieved_context"] = "죄송합니다만, 질문하신 내용은 PPT 슬라이드에 수록되어 있지 않은 내용입니다."
    else:
        state["retrieved_context"] = "\n---\n".join(context_docs)
        
    return state

def node_web_searcher(state: ChatState) -> ChatState:
    """[Web Searcher Worker] Tavily API를 사용해 PPT 범위 외부의 연관 정보를 실시간 구글링 수집합니다."""
    from langchain_tavily import TavilySearch
    query = state["user_query"]
    tavily_key = state["tavily_key"]
    
    try:
        tavily = TavilySearch(max_results=3, tavily_api_key=tavily_key)
        search_data = tavily.invoke(query)
        state["web_search_context"] = str(search_data)
    except Exception as e:
        state["web_search_context"] = f"(실시간 검색 오류: {str(e)})"
        
    return state

def node_response_generator(state: ChatState) -> ChatState:
    """[Response Generator Worker] 축적된 슬라이드 정보와 웹 검색 컨텍스트를 활용해 완성도 높은 강사 톤의 답변을 완성합니다."""
    query = state["user_query"]
    retrieved = state.get("retrieved_context", "")
    web_context = state.get("web_search_context", "")
    
    sys_instruction = """
    역할: 당신은 제공받은 PPT 슬라이드 지식 데이터베이스(Retrieved Context) 혹은 보충 검색 결과(Web Context)에 완벽하게 기반하여 답변하는 교육 비서입니다.
    
    규칙:
    - 반드시 주어진 컨텍스트 내에서 정확한 사실에 기반해 풍성하고 정중한 한글 구어체로 답변하세요.
    - 웹 검색 정보가 유입되었다면 최신 검색 내용을 바탕으로 친절하게 답변을 구성하세요.
    - 만약 내부 컨텍스트에 정보가 전혀 없고, 웹 검색(Tavily) 결과도 없는 상황이라면 임의로 창작하지 말고 다음 문장으로만 단정하게 대답을 시작하세요: "죄송합니다만, 질문하신 내용은 PPT 슬라이드 내에서 찾을 수 없었으며, 보충 검색 키가 없어 실시간 웹 서칭을 수행하지 못했습니다."
    """
    
    prompt_content = (
        f"[Retrieved Context]\n{retrieved}\n\n"
        f"[Web Context]\n{web_context if web_context else '(웹 검색 데이터 없음)'}\n\n"
        f"[사용자 질문]\n{query}"
    )
    
    chat = ChatOpenAI(model=LLM_MODEL, temperature=0.2, openai_api_key=state["openai_key"])
    response = chat.invoke([
        SystemMessage(content=sys_instruction),
        HumanMessage(content=prompt_content)
    ])
    
    state["final_response"] = response.content
    return state
